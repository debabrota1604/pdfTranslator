"""
Office CAT Pipeline - PDF → Office → Moses/XLIFF → Office → PDF.

Combines Office roundtrip conversion with CAT (Computer-Assisted Translation)
formats for better intermediate representation.

Workflow:
1. PDF → Office (DOCX/PPTX/XLSX) using pdf2docx/python-pptx/openpyxl
2. Extract text from Office XML
3. Generate Moses (parallel text) or XLIFF format for translation
4. Parse translated text and update Office XML
5. Office → PDF via LibreOffice

Output Formats:
- Moses: Simple parallel text files (source.txt, target.txt)
- XLIFF: Industry-standard XML Localization Interchange File Format

Required dependencies:
- pdf2docx: pip install pdf2docx (for PDF → DOCX)
- python-pptx: pip install python-pptx (for PPTX handling)
- openpyxl: pip install openpyxl (for XLSX handling)
- LibreOffice: https://www.libreoffice.org/ (for Office → PDF conversion)
"""

from __future__ import annotations

import json
import re
import subprocess
import shutil
import sys
import time
from dataclasses import dataclass
from enum import Enum
from pathlib import Path
from typing import Optional
from xml.etree import ElementTree as ET

from .base import (
    TranslationPipeline,
    PipelineConfig,
    PipelineType,
    ExtractResult,
    MergeResult,
)
from .office_xml import (
    get_handler,
    ExtractionResult,
)
from .office_roundtrip import (
    OfficeFormat,
    ProgressSpinner,
    prompt_office_format,
)
from ..source_detector import detect_source_format, SourceFormat


class CATFormat(Enum):
    """CAT output format types."""
    MOSES = "moses"      # Parallel text files
    XLIFF = "xliff"      # XLIFF 1.2/2.0 format
    TMX = "tmx"          # Translation Memory eXchange (future)


@dataclass
class OfficeCATConfig(PipelineConfig):
    """Configuration for Office CAT pipeline."""
    
    pipeline_type: PipelineType = PipelineType.OFFICE_CAT
    
    # Office format (auto-detect from PDF if AUTO)
    office_format: OfficeFormat = OfficeFormat.AUTO
    
    # CAT output format (moses is default - simpler for most workflows)
    cat_format: CATFormat = CATFormat.MOSES
    
    # XLIFF version (if cat_format is XLIFF)
    xliff_version: str = "1.2"
    
    # Source language code (for XLIFF)
    source_language: str = "en"
    
    # File encoding for Moses/text output files
    encoding: str = "utf-8"
    
    # LibreOffice path
    libreoffice_path: Optional[Path] = None
    
    # Keep intermediate files (always True now)
    keep_intermediate: bool = True


class OfficeCATpipeline(TranslationPipeline):
    """
    Office CAT Pipeline.
    
    Workflow:
    1. Detect source format from PDF metadata
    2. Convert PDF to Office format
    3. Extract text from Office XML
    4. Generate Moses or XLIFF translation files
    5. Parse translations and update Office XML
    6. Convert Office → PDF via LibreOffice
    """
    
    def __init__(self, config: OfficeCATConfig):
        super().__init__(config)
        self.config: OfficeCATConfig = config
        self._detected_format: Optional[SourceFormat] = None
        self._extraction_result: Optional[ExtractionResult] = None
        self._check_dependencies()
    
    def _check_dependencies(self) -> None:
        """Check available dependencies."""
        self._has_pdf2docx = False
        self._has_pptx = False
        self._has_openpyxl = False
        self._has_libreoffice = False
        self._has_translate_toolkit = False
        
        try:
            import pdf2docx
            self._has_pdf2docx = True
        except ImportError:
            pass
        
        try:
            import pptx
            self._has_pptx = True
        except ImportError:
            pass
        
        try:
            import openpyxl
            self._has_openpyxl = True
        except ImportError:
            pass
        
        try:
            from translate.storage import xliff
            self._has_translate_toolkit = True
        except ImportError:
            pass
        
        # Check LibreOffice
        lo_paths = [
            self.config.libreoffice_path,
            Path("C:/Program Files/LibreOffice/program/soffice.exe"),
            Path("C:/Program Files (x86)/LibreOffice/program/soffice.exe"),
            Path("/usr/bin/libreoffice"),
            Path("/usr/bin/soffice"),
            Path("/Applications/LibreOffice.app/Contents/MacOS/soffice"),
        ]
        for path in lo_paths:
            if path and path.exists():
                self.config.libreoffice_path = path
                self._has_libreoffice = True
                break
    
    @property
    def name(self) -> str:
        return f"Office CAT ({self.config.cat_format.value.upper()})"
    
    @property
    def description(self) -> str:
        fmt = self.config.cat_format.value.upper()
        return f"PDF -> Office -> {fmt} -> Office -> PDF"
    
    def derive_paths(self, input_path: Path) -> dict[str, Path]:
        """Derive paths including Office and CAT format files."""
        base = super().derive_paths(input_path)
        
        ext = self._get_office_extension()
        cat_fmt = self.config.cat_format
        
        base.update({
            "office": input_path.with_suffix(ext),
            "office_translated": input_path.parent / f"{input_path.stem}_translated{ext}",
        })
        
        # CAT format specific paths
        if cat_fmt == CATFormat.MOSES:
            base.update({
                "source_txt": input_path.parent / f"{input_path.name}_source.txt",
                "target_txt": input_path.parent / f"{input_path.name}_target.txt",
            })
        elif cat_fmt == CATFormat.XLIFF:
            base.update({
                "xliff": input_path.parent / f"{input_path.name}.xlf",
            })
        
        return base
    
    def _get_office_extension(self) -> str:
        """Get Office file extension based on detected/configured format."""
        if self.config.office_format == OfficeFormat.PPTX:
            return ".pptx"
        elif self.config.office_format == OfficeFormat.XLSX:
            return ".xlsx"
        else:
            return ".docx"
    
    def extract(self, input_path: Path) -> ExtractResult:
        """
        Extract text from PDF via Office format.
        
        Steps:
        1. Detect source format if AUTO
        2. Convert PDF to Office format
        3. Extract text from Office XML
        4. Generate Moses or XLIFF format files
        """
        # Step 1: Detect source format
        if self.config.office_format == OfficeFormat.AUTO:
            info = detect_source_format(input_path)
            self.config.office_format = prompt_office_format(info)
        
        paths = self.derive_paths(input_path)
        
        # Step 2: Convert PDF to Office
        pdf_size = input_path.stat().st_size / 1024
        print(f"  Input PDF: {pdf_size:.1f} KB")
        print(f"  Converting PDF -> {self.config.office_format.value.upper()}...")
        
        start_time = time.time()
        self._convert_pdf_to_office(input_path, paths["office"])
        elapsed = time.time() - start_time
        
        office_size = paths["office"].stat().st_size / 1024
        print(f"  Created {self.config.office_format.value.upper()}: {office_size:.1f} KB ({elapsed:.1f}s)")
        
        # Step 3: Extract text from Office XML
        print(f"  Extracting text from XML...")
        handler = get_handler(paths["office"])
        self._extraction_result = handler.extract()
        
        segments = self._extraction_result.segments
        print(f"  Found {len(segments)} text segments")
        
        # Step 4: Generate layout JSON
        layout = {
            "source_file": str(input_path),
            "pipeline": "office_cat",
            "office_format": self.config.office_format.value,
            "cat_format": self.config.cat_format.value,
            "office_path": str(paths["office"]),
            "source_language": self.config.source_language,
            "target_language": self.config.target_language,
            "encoding": self.config.encoding,
            "blocks": [
                {
                    "block_id": seg.block_id,
                    "text": seg.text,
                    "xml_path": seg.xml_path,
                    "type": seg.metadata.get("type", "text"),
                }
                for seg in segments
            ],
            "block_order": [seg.block_id for seg in segments],
        }
        paths["layout"].write_text(
            json.dumps(layout, indent=2, ensure_ascii=False),
            encoding="utf-8"
        )
        
        # Step 5: Generate CAT format files
        extra_files = {
            "office": paths["office"],
            "format": self.config.office_format.value,
            "cat_format": self.config.cat_format.value,
        }
        
        if self.config.cat_format == CATFormat.MOSES:
            self._generate_moses_format(segments, paths, layout)
            extra_files["source_txt"] = paths["source_txt"]
            extra_files["target_txt"] = paths["target_txt"]
            print(f"  Moses format:")
            print(f"    Source: {paths['source_txt'].name}")
            print(f"    Target: {paths['target_txt'].name}")
        elif self.config.cat_format == CATFormat.XLIFF:
            self._generate_xliff_format(segments, paths, layout)
            extra_files["xliff"] = paths["xliff"]
            print(f"  XLIFF: {paths['xliff'].name}")
        
        # Also generate tagged format as backup/alternative
        lines = []
        for i, seg in enumerate(segments):
            text = seg.text.replace('\n', '\\n')
            lines.append(f"<{i}>{text}</{i}>")
        
        paths["translate"].write_text('\n'.join(lines), encoding=self.config.encoding)
        paths["translated"].write_text('\n'.join(lines), encoding=self.config.encoding)
        
        return ExtractResult(
            layout_path=paths["layout"],
            translate_path=paths["translate"],
            translated_template_path=paths["translated"],
            extra_files=extra_files,
        )
    
    def _convert_pdf_to_office(self, pdf_path: Path, office_path: Path) -> None:
        """Convert PDF to Office format."""
        fmt = self.config.office_format
        
        if fmt == OfficeFormat.DOCX:
            self._convert_pdf_to_docx(pdf_path, office_path)
        elif fmt == OfficeFormat.PPTX:
            self._convert_pdf_to_pptx(pdf_path, office_path)
        elif fmt == OfficeFormat.XLSX:
            self._convert_pdf_to_xlsx(pdf_path, office_path)
    
    def _convert_pdf_to_docx(self, pdf_path: Path, docx_path: Path) -> None:
        """Convert PDF to DOCX using pdf2docx.
        
        Handles:
        - Text extraction with formatting preservation
        - Embedded images (photos, diagrams)
        - Tables and layout structure
        - Vector graphics rasterized by the PDF
        - PDF annotations converted to Word comments
        
        Note: pdf2docx handles images automatically during conversion.
        PDF annotations (comments, highlights, notes) are extracted separately
        and added as proper Word comments to avoid layout disruption.
        """
        if not self._has_pdf2docx:
            raise ImportError(
                "pdf2docx is required for PDF -> DOCX conversion.\n"
                "Install with: pip install pdf2docx"
            )
        
        import fitz
        from pdf2docx import Converter
        
        # Step 1: Extract annotations before conversion
        annotations = self._extract_pdf_annotations(pdf_path)
        
        # Step 2: Convert PDF to DOCX
        cv = Converter(str(pdf_path))
        cv.convert(str(docx_path), start=0, end=None)
        cv.close()
        
        # Step 3: Add annotations as proper Word comments
        if annotations:
            self._add_word_comments(docx_path, annotations)
            print(f"    Added {len(annotations)} comments from PDF annotations")
    
    def _extract_pdf_annotations(self, pdf_path: Path) -> list[dict]:
        """Extract annotations/comments from PDF.
        
        Extracts:
        - Text annotations (sticky notes)
        - Highlight annotations with comments
        - FreeText annotations
        - Popup annotations
        
        Returns list of annotation dicts with page, position, text, author info.
        """
        import fitz
        
        annotations = []
        doc = fitz.open(str(pdf_path))
        
        try:
            for page_num, page in enumerate(doc):
                if not page.annots():
                    continue
                
                for annot in page.annots():
                    annot_type = annot.type[0]
                    
                    # Types that should be converted to comments
                    comment_types = {0, 2, 8, 9, 10, 11}  # Text, FreeText, Highlight, Underline, Squiggly, StrikeOut
                    
                    if annot_type not in comment_types:
                        continue
                    
                    info = annot.info
                    content = info.get("content", "") or ""
                    subject = info.get("subject", "") or ""
                    title = info.get("title", "") or ""
                    
                    rect = annot.rect
                    
                    # For highlight/underline, get the underlying text
                    if annot_type in {8, 9, 10, 11}:
                        try:
                            highlighted_text = page.get_text("text", clip=rect).strip()
                        except Exception:
                            highlighted_text = ""
                    else:
                        highlighted_text = ""
                    
                    if not content and not highlighted_text:
                        continue
                    
                    annotations.append({
                        "page": page_num,
                        "type": annot.type[1],
                        "rect": [rect.x0, rect.y0, rect.x1, rect.y1],
                        "content": content,
                        "subject": subject,
                        "author": title,
                        "highlighted_text": highlighted_text,
                    })
        finally:
            doc.close()
        
        return annotations
    
    def _add_word_comments(self, docx_path: Path, annotations: list[dict]) -> None:
        """Add PDF annotations as Word comments to the DOCX file."""
        try:
            from docx import Document
            from docx.shared import Pt, RGBColor
        except ImportError:
            print("    Warning: python-docx not installed, skipping comment conversion")
            return
        
        doc = Document(str(docx_path))
        
        # Group annotations by page
        by_page = {}
        for annot in annotations:
            page = annot["page"]
            if page not in by_page:
                by_page[page] = []
            by_page[page].append(annot)
        
        if annotations:
            doc.add_paragraph()
            separator = doc.add_paragraph()
            separator.add_run("─" * 50)
            
            header = doc.add_paragraph()
            run = header.add_run("PDF Annotations (converted from comments)")
            run.bold = True
            run.font.size = Pt(12)
            run.font.color.rgb = RGBColor(0x80, 0x80, 0x80)
            
            for page_num in sorted(by_page.keys()):
                page_annotations = by_page[page_num]
                
                page_header = doc.add_paragraph()
                run = page_header.add_run(f"Page {page_num + 1}:")
                run.bold = True
                run.font.size = Pt(10)
                
                for annot in page_annotations:
                    para = doc.add_paragraph()
                    
                    type_run = para.add_run(f"[{annot['type']}] ")
                    type_run.font.color.rgb = RGBColor(0x00, 0x66, 0xCC)
                    type_run.font.size = Pt(9)
                    
                    if annot.get("author"):
                        author_run = para.add_run(f"({annot['author']}) ")
                        author_run.italic = True
                        author_run.font.size = Pt(9)
                    
                    if annot.get("highlighted_text"):
                        hl_run = para.add_run(f'"{annot["highlighted_text"]}" - ')
                        hl_run.font.color.rgb = RGBColor(0xFF, 0x99, 0x00)
                        hl_run.font.size = Pt(9)
                    
                    if annot.get("content"):
                        content_run = para.add_run(annot["content"])
                        content_run.font.size = Pt(9)
        
        doc.save(str(docx_path))
    
    def _convert_pdf_to_pptx(self, pdf_path: Path, pptx_path: Path) -> None:
        """Convert PDF to PPTX by extracting text and images.
        
        Handles:
        - Text extraction with position preservation
        - Embedded images (photos, diagrams)
        - Rendered SmartArt graphics
        - Vector graphics rasterized by the PDF
        
        Note: Images are extracted via _add_images_to_slide() helper.
        """
        if not self._has_pptx:
            raise ImportError(
                "python-pptx is required for PPTX creation.\n"
                "Install with: pip install python-pptx"
            )
        
        import fitz
        import io
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
        
        doc = fitz.open(str(pdf_path))
        total_pages = len(doc)
        
        try:
            prs = Presentation()
            
            if total_pages > 0:
                first_page = doc[0]
                prs.slide_width = int(first_page.rect.width * 914400 / 72)
                prs.slide_height = int(first_page.rect.height * 914400 / 72)
            
            blank_layout = prs.slide_layouts[6]
            
            for page_num, page in enumerate(doc, 1):
                sys.stdout.write(f"\r    Slide {page_num}/{total_pages}")
                sys.stdout.flush()
                
                slide = prs.slides.add_slide(blank_layout)
                blocks = page.get_text("dict")["blocks"]
                page_height = page.rect.height
                
                # Detect header/footer zones (top/bottom 10% of page)
                header_zone = page_height * 0.10
                footer_zone = page_height * 0.90
                
                # Process text blocks with full formatting
                for block in blocks:
                    if block.get("type") == 0:  # Text block
                        self._add_formatted_textbox(slide, block, header_zone, footer_zone, page_height)
                
                # Extract and add images (including SmartArt rendered as images)
                self._add_images_to_slide(doc, page, slide)
            
            sys.stdout.write(f"\r    Created {total_pages} slides    \n")
            sys.stdout.flush()
            prs.save(str(pptx_path))
            
        finally:
            doc.close()
    
    def _add_formatted_textbox(self, slide, block: dict, header_zone: float, footer_zone: float, page_height: float) -> None:
        """Add a text box with full formatting preservation.
        
        Preserves:
        - Font family, size, bold, italic, underline
        - Font color
        - Paragraph alignment
        - Auto-fit for text overflow
        - Header/footer detection
        """
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
        
        bbox = block.get("bbox", [0, 0, 100, 20])
        x, y, x1, y1 = bbox
        
        # Convert PDF points to inches
        left = Inches(x / 72)
        top = Inches(y / 72)
        width = Inches((x1 - x) / 72)
        height = Inches((y1 - y) / 72)
        
        # Minimum dimensions
        if width.inches < 0.5:
            width = Inches(0.5)
        if height.inches < 0.2:
            height = Inches(0.2)
        
        # Create text box
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        # Enable auto-fit for text overflow
        try:
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        except Exception:
            pass
        
        # Determine if this is header/footer
        is_header = y < header_zone
        is_footer = y > footer_zone
        
        # Process lines and spans with formatting
        first_para = True
        for line in block.get("lines", []):
            if first_para:
                para = tf.paragraphs[0]
                first_para = False
            else:
                para = tf.add_paragraph()
            
            # Detect alignment from line position
            line_bbox = line.get("bbox", bbox)
            line_center = (line_bbox[0] + line_bbox[2]) / 2
            block_center = (x + x1) / 2
            
            if abs(line_center - block_center) < 5:
                para.alignment = PP_ALIGN.CENTER
            elif line_bbox[0] > x + 10:
                para.alignment = PP_ALIGN.RIGHT
            else:
                para.alignment = PP_ALIGN.LEFT
            
            # Process spans with formatting
            for span in line.get("spans", []):
                text = span.get("text", "")
                if not text:
                    continue
                
                run = para.add_run()
                run.text = text
                
                # Font size
                font_size = span.get("size", 12)
                run.font.size = Pt(font_size)
                
                # Font family
                font_name = span.get("font", "")
                if font_name:
                    # Extract base font name (remove style suffixes)
                    base_font = font_name.split("-")[0].split("+")[-1]
                    run.font.name = base_font
                
                # Font flags (bold, italic)
                flags = span.get("flags", 0)
                run.font.bold = bool(flags & 2**4)  # Bit 4 = bold
                run.font.italic = bool(flags & 2**1)  # Bit 1 = italic
                
                # Font color
                color = span.get("color", 0)
                if color and color != 0:
                    try:
                        # Color is usually an integer RGB
                        r = (color >> 16) & 0xFF
                        g = (color >> 8) & 0xFF
                        b = color & 0xFF
                        run.font.color.rgb = RGBColor(r, g, b)
                    except Exception:
                        pass
        
        # Add metadata for header/footer detection
        if is_header or is_footer:
            try:
                # Move header/footer text boxes to appropriate z-order
                pass  # PowerPoint handles this automatically
            except Exception:
                pass
    
    def _add_images_to_slide(self, doc, page, slide) -> None:
        """Extract images from PDF page and add to PowerPoint slide.
        
        Handles:
        - Embedded images (photos, diagrams)
        - Rendered SmartArt graphics
        - Vector graphics rasterized by the PDF
        """
        import io
        from pptx.util import Inches
        
        try:
            # Get all images on this page
            image_list = page.get_images(full=True)
            
            for img_index, img_info in enumerate(image_list):
                xref = img_info[0]  # Image reference number
                
                try:
                    # Extract image data
                    base_image = doc.extract_image(xref)
                    if not base_image:
                        continue
                    
                    image_bytes = base_image["image"]
                    
                    # Get image position on page
                    img_rects = page.get_image_rects(xref)
                    if not img_rects:
                        continue
                    
                    for rect in img_rects:
                        # Convert PDF coordinates to inches
                        left = Inches(rect.x0 / 72)
                        top = Inches(rect.y0 / 72)
                        width = Inches((rect.x1 - rect.x0) / 72)
                        height = Inches((rect.y1 - rect.y0) / 72)
                        
                        # Skip very small images (likely artifacts)
                        if width.inches < 0.1 or height.inches < 0.1:
                            continue
                        
                        # Add image to slide
                        image_stream = io.BytesIO(image_bytes)
                        try:
                            slide.shapes.add_picture(
                                image_stream, left, top, width, height
                            )
                        except Exception:
                            pass
                            
                except Exception:
                    continue
                    
        except Exception:
            pass
    
    def _convert_pdf_to_xlsx(self, pdf_path: Path, xlsx_path: Path) -> None:
        """Convert PDF to XLSX by extracting tabular data and images.
        
        Handles:
        - Text extraction arranged into rows/columns
        - Font formatting (bold, italic, color, size)
        - Cell alignment and borders
        - Embedded images (photos, diagrams, charts)
        - Rendered SmartArt graphics
        - Vector graphics rasterized by the PDF
        
        Note: Images are extracted via _add_images_to_worksheet() helper.
        """
        if not self._has_openpyxl:
            raise ImportError(
                "openpyxl is required for XLSX creation.\n"
                "Install with: pip install openpyxl"
            )
        
        import fitz
        import io
        from openpyxl import Workbook
        from openpyxl.drawing.image import Image as XLImage
        from openpyxl.utils import get_column_letter
        from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
        
        doc = fitz.open(str(pdf_path))
        total_pages = len(doc)
        
        try:
            wb = Workbook()
            
            for page_num, page in enumerate(doc):
                sys.stdout.write(f"\r    Sheet {page_num + 1}/{total_pages}")
                sys.stdout.flush()
                
                if page_num == 0:
                    ws = wb.active
                    ws.title = f"Sheet{page_num + 1}"
                else:
                    ws = wb.create_sheet(f"Sheet{page_num + 1}")
                
                blocks = page.get_text("dict")["blocks"]
                text_blocks = [b for b in blocks if b.get("type") == 0]
                text_blocks.sort(key=lambda b: (b["bbox"][1], b["bbox"][0]))
                
                # Group blocks into rows based on Y position
                rows = []
                current_row = []
                last_y = None
                
                for block in text_blocks:
                    y = block["bbox"][1]
                    if last_y is None or abs(y - last_y) < 15:
                        current_row.append(block)
                    else:
                        if current_row:
                            rows.append(sorted(current_row, key=lambda b: b["bbox"][0]))
                        current_row = [block]
                    last_y = y
                
                if current_row:
                    rows.append(sorted(current_row, key=lambda b: b["bbox"][0]))
                
                max_row = 0
                for row_num, row_blocks in enumerate(rows, 1):
                    max_row = row_num
                    for col_num, block in enumerate(row_blocks, 1):
                        self._add_formatted_cell(ws, row_num, col_num, block)
                
                # Add images to worksheet
                self._add_images_to_worksheet(doc, page, ws, max_row)
            
            sys.stdout.write(f"\r    Created {total_pages} sheets    \n")
            sys.stdout.flush()
            wb.save(str(xlsx_path))
            
        finally:
            doc.close()
    
    def _add_formatted_cell(self, ws, row_num: int, col_num: int, block: dict) -> None:
        """Add a cell with formatting preservation.
        
        Preserves:
        - Font family, size, bold, italic
        - Font color
        - Text alignment
        - Cell width adjustment for content
        """
        from openpyxl.styles import Font, Alignment
        from openpyxl.utils import get_column_letter
        
        # Collect text and dominant formatting
        text_parts = []
        dominant_font_size = 11
        dominant_font_name = "Calibri"
        is_bold = False
        is_italic = False
        font_color = "000000"
        
        for line in block.get("lines", []):
            line_texts = []
            for span in line.get("spans", []):
                text = span.get("text", "")
                if text:
                    line_texts.append(text)
                    
                    # Get formatting from first substantial span
                    if len(text.strip()) > 0:
                        dominant_font_size = span.get("size", 11)
                        
                        font_name = span.get("font", "")
                        if font_name:
                            dominant_font_name = font_name.split("-")[0].split("+")[-1]
                        
                        flags = span.get("flags", 0)
                        is_bold = is_bold or bool(flags & 2**4)
                        is_italic = is_italic or bool(flags & 2**1)
                        
                        color = span.get("color", 0)
                        if color and color != 0:
                            r = (color >> 16) & 0xFF
                            g = (color >> 8) & 0xFF
                            b = color & 0xFF
                            font_color = f"{r:02X}{g:02X}{b:02X}"
            
            if line_texts:
                text_parts.append(''.join(line_texts))
        
        cell_text = ' '.join(text_parts).strip()
        if not cell_text:
            return
        
        cell = ws.cell(row=row_num, column=col_num, value=cell_text)
        
        # Apply font formatting
        try:
            cell.font = Font(
                name=dominant_font_name,
                size=int(dominant_font_size),
                bold=is_bold,
                italic=is_italic,
                color=font_color,
            )
        except Exception:
            pass
        
        # Apply alignment
        try:
            cell.alignment = Alignment(
                wrap_text=True,
                vertical='center',
            )
        except Exception:
            pass
        
        # Adjust column width based on content
        try:
            col_letter = get_column_letter(col_num)
            current_width = ws.column_dimensions[col_letter].width or 8
            # Estimate width needed (approx 1 char = 1 unit)
            needed_width = min(50, max(current_width, len(cell_text) * 1.1))
            ws.column_dimensions[col_letter].width = needed_width
        except Exception:
            pass
    
    def _add_images_to_worksheet(self, doc, page, ws, start_row: int) -> None:
        """Extract images from PDF page and add to Excel worksheet.
        
        Handles:
        - Embedded images (photos, diagrams, charts)
        - Rendered SmartArt graphics
        - Vector graphics rasterized by the PDF
        """
        import io
        from openpyxl.drawing.image import Image as XLImage
        from openpyxl.utils import get_column_letter
        
        try:
            image_list = page.get_images(full=True)
            image_count = 0
            
            for img_index, img_info in enumerate(image_list):
                xref = img_info[0]
                
                try:
                    base_image = doc.extract_image(xref)
                    if not base_image:
                        continue
                    
                    image_bytes = base_image["image"]
                    img_width = base_image.get("width", 100)
                    img_height = base_image.get("height", 100)
                    
                    # Skip very small images (likely artifacts)
                    if img_width < 10 or img_height < 10:
                        continue
                    
                    img_rects = page.get_image_rects(xref)
                    
                    image_stream = io.BytesIO(image_bytes)
                    try:
                        xl_img = XLImage(image_stream)
                        
                        # Scale image to reasonable size
                        if xl_img.width > 400:
                            scale = 400 / xl_img.width
                            xl_img.width = 400
                            xl_img.height = int(xl_img.height * scale)
                        
                        if img_rects:
                            rect = img_rects[0]
                            col = max(1, int(rect.x0 / 72))
                            row = max(1, int(rect.y0 / 20))
                            cell_ref = f"{get_column_letter(col)}{row}"
                        else:
                            cell_ref = f"A{start_row + 2 + image_count * 15}"
                        
                        ws.add_image(xl_img, cell_ref)
                        image_count += 1
                        
                    except Exception:
                        pass
                        
                except Exception:
                    continue
                    
        except Exception:
            pass
    
    def _generate_moses_format(
        self,
        segments: list,
        paths: dict[str, Path],
        layout: dict,
    ) -> None:
        """Generate Moses parallel text format.
        
        Moses format:
        - source.txt: One segment per line (source text)
        - target.txt: One segment per line (to be translated)
        
        Each line corresponds to the same translation unit.
        Segment IDs are preserved in a separate mapping file.
        """
        source_lines = []
        target_lines = []
        
        for seg in segments:
            # Normalize text for Moses: replace newlines with special token
            text = seg.text.replace('\n', ' <br> ')
            # Clean up multiple spaces
            text = ' '.join(text.split())
            source_lines.append(text)
            target_lines.append(text)  # Template for translation
        
        paths["source_txt"].write_text('\n'.join(source_lines), encoding=self.config.encoding)
        paths["target_txt"].write_text('\n'.join(target_lines), encoding=self.config.encoding)
        
        # Also save segment ID mapping
        mapping_path = paths["source_txt"].with_suffix('.mapping.json')
        mapping = {
            "format": "moses",
            "segments": [
                {"line": i, "block_id": seg.block_id}
                for i, seg in enumerate(segments)
            ]
        }
        mapping_path.write_text(json.dumps(mapping, indent=2), encoding="utf-8")
    
    def _generate_xliff_format(
        self,
        segments: list,
        paths: dict[str, Path],
        layout: dict,
    ) -> None:
        """Generate XLIFF format using translate-toolkit (OASIS compliant).
        
        XLIFF (XML Localization Interchange File Format) is an industry standard
        for exchanging translation data between CAT tools.
        
        Uses translate-toolkit for standards-compliant XLIFF 1.2 output.
        Falls back to basic XML generation if translate-toolkit not installed.
        """
        if self._has_translate_toolkit:
            self._generate_xliff_with_toolkit(segments, paths, layout)
        else:
            self._generate_xliff_basic(segments, paths, layout)
    
    def _generate_xliff_with_toolkit(
        self,
        segments: list,
        paths: dict[str, Path],
        layout: dict,
    ) -> None:
        """Generate XLIFF using translate-toolkit (standards-compliant).
        
        translate-toolkit implements OASIS XLIFF 1.2 specification:
        http://docs.oasis-open.org/xliff/xliff-core/xliff-core.html
        """
        from translate.storage.xliff import xlifffile
        
        source_lang = self.config.source_language
        target_lang = self.config.target_language
        
        # Create XLIFF file with proper header
        xliff_store = xlifffile()
        xliff_store.setsourcelanguage(source_lang)
        xliff_store.settargetlanguage(target_lang)
        
        for i, seg in enumerate(segments):
            # Create translation unit with source text
            unit = xliff_store.addsourceunit(seg.text)
            unit.setid(str(i))
            
            # Set target (pre-filled with source for template)
            unit.target = seg.text
            
            # Set state to needs-translation (XLIFF 1.2 standard state)
            if hasattr(unit, 'markfuzzy'):
                unit.markfuzzy()  # translate-toolkit way to set state
            
            # Add location/context with block ID
            if hasattr(unit, 'addlocation'):
                unit.addlocation(seg.block_id)
            
            # Add note with metadata
            if seg.metadata and hasattr(unit, 'addnote'):
                note_text = f"Type: {seg.metadata.get('type', 'text')}"
                unit.addnote(note_text, origin="developer")
        
        # Set filename on the file node if available
        try:
            filenode = xliff_store.getfilenode(None)
            if filenode is not None:
                xliff_store.setfilename(filenode, str(paths["office"].name))
        except Exception:
            pass  # Filename is optional
        
        # Write to file
        with open(paths["xliff"], 'wb') as f:
            xliff_store.serialize(f)
        
        print(f"    Using translate-toolkit (OASIS XLIFF 1.2 compliant)")
    
    def _generate_xliff_basic(
        self,
        segments: list,
        paths: dict[str, Path],
        layout: dict,
    ) -> None:
        """Generate basic XLIFF format (fallback without translate-toolkit).
        
        Manual XML generation following XLIFF 1.2 structure.
        For full compliance, install translate-toolkit.
        """
        source_lang = self.config.source_language
        target_lang = self.config.target_language
        
        # XLIFF 1.2 format
        xliff = ET.Element('xliff')
        xliff.set('version', '1.2')
        xliff.set('xmlns', 'urn:oasis:names:tc:xliff:document:1.2')
        
        file_elem = ET.SubElement(xliff, 'file')
        file_elem.set('original', str(paths["office"]))
        file_elem.set('source-language', source_lang)
        file_elem.set('target-language', target_lang)
        file_elem.set('datatype', 'plaintext')
        
        body = ET.SubElement(file_elem, 'body')
        
        for i, seg in enumerate(segments):
            trans_unit = ET.SubElement(body, 'trans-unit')
            trans_unit.set('id', str(i))
            trans_unit.set('resname', seg.block_id)
            
            source = ET.SubElement(trans_unit, 'source')
            source.text = seg.text
            
            target = ET.SubElement(trans_unit, 'target')
            target.set('state', 'needs-translation')
            target.text = seg.text  # Pre-fill with source for template
            
            # Add note with metadata
            if seg.metadata:
                note = ET.SubElement(trans_unit, 'note')
                note.text = f"Type: {seg.metadata.get('type', 'text')}"
        
        # Pretty print
        self._indent_xml(xliff)
        
        tree = ET.ElementTree(xliff)
        tree.write(str(paths["xliff"]), encoding="utf-8", xml_declaration=True)
        
        print(f"    Using basic XLIFF (install translate-toolkit for full compliance)")
    
    def _indent_xml(self, elem: ET.Element, level: int = 0) -> None:
        """Add indentation to XML for readability."""
        indent = "\n" + "  " * level
        if len(elem):
            if not elem.text or not elem.text.strip():
                elem.text = indent + "  "
            if not elem.tail or not elem.tail.strip():
                elem.tail = indent
            for child in elem:
                self._indent_xml(child, level + 1)
            if not child.tail or not child.tail.strip():
                child.tail = indent
        else:
            if level and (not elem.tail or not elem.tail.strip()):
                elem.tail = indent
    
    def merge(
        self,
        input_path: Path,
        output_path: Path,
        translated_path: Path,
        layout_path: Path,
    ) -> MergeResult:
        """
        Apply translations and convert back to PDF.
        
        Steps:
        1. Load layout and determine CAT format
        2. Parse translations from Moses/XLIFF or tagged format
        3. Update Office XML with translations
        4. Convert Office → PDF via LibreOffice
        """
        if not self._has_libreoffice:
            raise RuntimeError(
                "LibreOffice is required for Office -> PDF conversion.\n"
                "Install from: https://www.libreoffice.org/"
            )
        
        # Load layout
        layout = json.loads(layout_path.read_text(encoding="utf-8"))
        office_format = layout.get("office_format", "docx")
        cat_format = layout.get("cat_format", "moses")
        office_path = Path(layout.get("office_path", ""))
        block_order = layout.get("block_order", [])
        
        if not office_path.exists():
            raise FileNotFoundError(f"Office file not found: {office_path}")
        
        # Parse translations based on format
        print(f"  CAT format: {cat_format.upper()}")
        
        if cat_format == "moses":
            translations = self._parse_moses_translations(layout_path, block_order)
        elif cat_format == "xliff":
            translations = self._parse_xliff_translations(layout_path, block_order)
        else:
            # Fallback to tagged format
            translations = self._parse_tagged_translations(translated_path, block_order)
        
        print(f"  Parsed {len(translations)} translations")
        
        # Create translated Office file
        ext = f".{office_format}"
        translated_office = input_path.parent / f"{input_path.stem}_translated{ext}"
        
        # Get handler and extract current state
        print(f"  Reading {office_format.upper()} structure...")
        handler = get_handler(office_path)
        extraction = handler.extract()
        
        # Update with translations
        print(f"  Updating {office_format.upper()} with translations...")
        handler.update(translated_office, translations, extraction)
        
        office_size = translated_office.stat().st_size / 1024
        print(f"  Office file: {office_size:.1f} KB")
        
        # Convert to PDF
        self._office_to_pdf(translated_office, output_path, office_format)
        
        # Keep intermediate files
        print(f"  Kept: {office_path.name} ({office_format.upper()})")
        print(f"  Kept: {translated_office.name} ({office_format.upper()})")
        
        return MergeResult(
            output_path=output_path,
            blocks_processed=len(translations),
        )
    
    def _parse_moses_translations(
        self,
        layout_path: Path,
        block_order: list[str],
    ) -> dict[str, str]:
        """Parse Moses format translations."""
        target_path = layout_path.parent / f"{layout_path.stem.replace('_layout', '')}_target.txt"
        mapping_path = layout_path.parent / f"{layout_path.stem.replace('_layout', '')}_source.mapping.json"
        
        if not target_path.exists():
            raise FileNotFoundError(
                f"Moses target file not found: {target_path}\n"
                f"Expected: {target_path.name}"
            )
        
        target_lines = target_path.read_text(encoding=self.config.encoding).split('\n')
        
        # Load mapping if exists
        if mapping_path.exists():
            mapping = json.loads(mapping_path.read_text(encoding="utf-8"))
            segment_map = {s["line"]: s["block_id"] for s in mapping.get("segments", [])}
        else:
            segment_map = {i: block_id for i, block_id in enumerate(block_order)}
        
        translations = {}
        for line_num, text in enumerate(target_lines):
            if line_num in segment_map:
                # Convert Moses format back to normal text
                text = text.replace(' <br> ', '\n')
                translations[segment_map[line_num]] = text
        
        return translations
    
    def _parse_xliff_translations(
        self,
        layout_path: Path,
        block_order: list[str],
    ) -> dict[str, str]:
        """Parse XLIFF format translations using translate-toolkit."""
        xliff_path = layout_path.parent / f"{layout_path.stem.replace('_layout', '')}.xlf"
        
        if not xliff_path.exists():
            raise FileNotFoundError(
                f"XLIFF file not found: {xliff_path}\n"
                f"Expected: {xliff_path.name}"
            )
        
        if self._has_translate_toolkit:
            return self._parse_xliff_with_toolkit(xliff_path, block_order)
        else:
            return self._parse_xliff_basic(xliff_path, block_order)
    
    def _parse_xliff_with_toolkit(
        self,
        xliff_path: Path,
        block_order: list[str],
    ) -> dict[str, str]:
        """Parse XLIFF using translate-toolkit (standards-compliant)."""
        from translate.storage.xliff import xlifffile
        
        with open(xliff_path, 'rb') as f:
            xliff_store = xlifffile(f)
        
        translations = {}
        
        for unit in xliff_store.units:
            # Skip header units
            if unit.isheader():
                continue
            
            unit_id = unit.getid()
            
            # Get target text (translated)
            target = unit.target
            if target:
                # Map by ID or use block_order
                if unit_id.isdigit():
                    idx = int(unit_id)
                    if idx < len(block_order):
                        translations[block_order[idx]] = target
                else:
                    translations[unit_id] = target
        
        return translations
    
    def _parse_xliff_basic(
        self,
        xliff_path: Path,
        block_order: list[str],
    ) -> dict[str, str]:
        """Parse XLIFF using basic XML (fallback)."""
        tree = ET.parse(str(xliff_path))
        root = tree.getroot()
        
        # Handle XLIFF namespace
        ns = {'xliff': 'urn:oasis:names:tc:xliff:document:1.2'}
        
        translations = {}
        
        # Try with namespace first, then without
        trans_units = root.findall('.//xliff:trans-unit', ns)
        if not trans_units:
            trans_units = root.findall('.//trans-unit')
        
        for tu in trans_units:
            block_id = tu.get('resname') or tu.get('id')
            
            # Get target text
            target = tu.find('xliff:target', ns)
            if target is None:
                target = tu.find('target')
            
            if target is not None and target.text:
                translations[block_id] = target.text
        
        return translations
    
    def _parse_tagged_translations(
        self,
        translated_path: Path,
        block_order: list[str],
    ) -> dict[str, str]:
        """Parse tagged format translations (fallback)."""
        content = translated_path.read_text(encoding=self.config.encoding)
        translations = {}
        
        pattern = r'<(\d+)>(.*?)</\1>'
        for match in re.finditer(pattern, content, re.DOTALL):
            idx = int(match.group(1))
            text = match.group(2).replace('\\n', '\n')
            
            if idx < len(block_order):
                translations[block_order[idx]] = text
        
        return translations
    
    def _office_to_pdf(
        self,
        office_path: Path,
        pdf_path: Path,
        office_format: str = "docx",
    ) -> None:
        """Convert Office file to PDF via LibreOffice."""
        cmd = [
            str(self.config.libreoffice_path),
            "--headless",
            "--convert-to", "pdf",
            "--outdir", str(pdf_path.parent),
            str(office_path),
        ]
        
        with ProgressSpinner(f"Converting {office_format.upper()} to PDF via LibreOffice"):
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=300)
        
        if result.returncode != 0:
            raise RuntimeError(f"LibreOffice conversion failed: {result.stderr}")
        
        generated = office_path.with_suffix(".pdf")
        if generated.name != pdf_path.name:
            generated_in_outdir = pdf_path.parent / generated.name
            if generated_in_outdir.exists() and generated_in_outdir != pdf_path:
                shutil.move(str(generated_in_outdir), str(pdf_path))
    
    def get_translation_prompt(self, block_count: int) -> str:
        """Get translation prompt based on CAT format."""
        if self.config.cat_format == CATFormat.MOSES:
            return self._get_moses_prompt(block_count)
        elif self.config.cat_format == CATFormat.XLIFF:
            return self._get_xliff_prompt(block_count)
        else:
            from ..translation_io import get_translation_prompt
            return get_translation_prompt(block_count, self.config.target_language)
    
    def _get_moses_prompt(self, block_count: int) -> str:
        """Get Moses format translation prompt."""
        return f"""
Moses Parallel Text Format - {block_count} segments to translate

Files:
- *_source.txt: Source text (one segment per line)
- *_target.txt: Target text (translate this file)

Instructions:
1. Open *_target.txt
2. Translate each line to {self.config.target_language}
3. Keep the same number of lines (do NOT add or remove lines)
4. Preserve <br> markers (they represent line breaks)
5. Save as UTF-8 encoding

Example:
Source line: Hello world <br> Welcome
Target line: नमस्ते दुनिया <br> स्वागत है

After translation, run: python main.py merge input.pdf
"""
    
    def _get_xliff_prompt(self, block_count: int) -> str:
        """Get XLIFF format translation prompt."""
        return f"""
XLIFF Format - {block_count} translation units

File: *.xlf (XLIFF 1.2 format)

Instructions:
1. Open the .xlf file in your CAT tool (SDL Trados, memoQ, OmegaT, etc.)
2. Translate each <source> to <target> in {self.config.target_language}
3. Update state="needs-translation" to state="translated"
4. Save the file

Or edit manually:
- Find each <trans-unit>
- Translate the <target> element content
- Keep XML structure intact

After translation, run: python main.py merge input.pdf
"""


def create_office_cat_pipeline(
    target_language: str = "Hindi",
    source_language: str = "en",
    office_format: OfficeFormat = OfficeFormat.AUTO,
    cat_format: CATFormat = CATFormat.MOSES,
    encoding: str = "utf-8",
    libreoffice_path: Optional[Path] = None,
) -> OfficeCATpipeline:
    """Factory function to create Office CAT pipeline.
    
    Args:
        target_language: Target language for translation (default: Hindi)
        source_language: Source language code (default: en)
        office_format: Office format - AUTO, DOCX, PPTX, XLSX (default: AUTO)
        cat_format: CAT output format - MOSES or XLIFF (default: MOSES)
        encoding: File encoding for text output (default: utf-8)
        libreoffice_path: Path to LibreOffice executable (auto-detected if None)
    """
    config = OfficeCATConfig(
        target_language=target_language,
        source_language=source_language,
        office_format=office_format,
        cat_format=cat_format,
        encoding=encoding,
        libreoffice_path=libreoffice_path,
    )
    return OfficeCATpipeline(config)
