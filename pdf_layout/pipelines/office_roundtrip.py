"""
Office Roundtrip Pipeline - PDF → Office Format → Translate → Office → PDF.

Workflow:
1. PDF → Office (DOCX/PPTX/XLSX)
   - DOCX: pdf2docx library (high fidelity conversion)
   - PPTX/XLSX: Create from PDF text extraction
   
2. Extract XML from Office (ZIP with XML inside)
   - Uses OfficeXMLHandler to parse document.xml, slides, sheets
   
3. Generate translation file with tagged text
   
4. Apply translations to XML and repack Office file

5. Office → PDF via LibreOffice

Supports:
- DOCX (Word documents)
- PPTX (PowerPoint presentations)  
- XLSX (Excel spreadsheets)

Required dependencies (install separately):
- pdf2docx: pip install pdf2docx (for PDF → DOCX)
- python-pptx: pip install python-pptx (for PPTX handling)
- openpyxl: pip install openpyxl (for XLSX handling)
- LibreOffice: https://www.libreoffice.org/ (for Office → PDF conversion)
"""

from __future__ import annotations

import json
import re
import subprocess
import shutil
import sys
import threading
import time
from dataclasses import dataclass
from enum import Enum
from pathlib import Path
from typing import Optional

from .base import (
    TranslationPipeline,
    PipelineConfig,
    PipelineType,
    ExtractResult,
    MergeResult,
)
from .office_xml import (
    OfficeXMLHandler,
    DocxXMLHandler,
    PptxXMLHandler,
    XlsxXMLHandler,
    ExtractionResult,
    get_handler,
)
from ..source_detector import detect_source_format, SourceFormat, SourceInfo


class OfficeFormat(Enum):
    """Office format types."""
    DOCX = "docx"
    PPTX = "pptx"
    XLSX = "xlsx"
    AUTO = "auto"  # Auto-detect from PDF metadata


class ProgressSpinner:
    """
    Threaded progress spinner for long-running operations.
    
    Usage:
        with ProgressSpinner("Converting..."):
            do_long_operation()
    """
    
    # ASCII frames for Windows compatibility
    FRAMES = ["-", "\\", "|", "/"]
    
    def __init__(self, message: str = "Processing"):
        self.message = message
        self._stop_event = threading.Event()
        self._thread: Optional[threading.Thread] = None
        self._start_time = 0.0
    
    def _spin(self) -> None:
        """Spinner thread function."""
        frame_idx = 0
        while not self._stop_event.is_set():
            elapsed = time.time() - self._start_time
            frame = self.FRAMES[frame_idx % len(self.FRAMES)]
            sys.stdout.write(f"\r  [{frame}] {self.message} ({elapsed:.1f}s)")
            sys.stdout.flush()
            frame_idx += 1
            time.sleep(0.2)
    
    def __enter__(self) -> "ProgressSpinner":
        self._start_time = time.time()
        self._stop_event.clear()
        self._thread = threading.Thread(target=self._spin, daemon=True)
        self._thread.start()
        return self
    
    def __exit__(self, exc_type, exc_val, exc_tb) -> None:
        self._stop_event.set()
        if self._thread:
            self._thread.join(timeout=1.0)
        elapsed = time.time() - self._start_time
        # Clear spinner line and write completion
        sys.stdout.write(f"\r  [OK] {self.message} ({elapsed:.1f}s)    \n")
        sys.stdout.flush()


@dataclass
class OfficeRoundtripConfig(PipelineConfig):
    """Configuration for Office roundtrip pipeline."""
    
    pipeline_type: PipelineType = PipelineType.DOCX_ROUNDTRIP
    
    # Office format (auto-detect from PDF if AUTO)
    office_format: OfficeFormat = OfficeFormat.AUTO
    
    # LibreOffice path
    libreoffice_path: Optional[Path] = None
    
    # Keep intermediate files
    keep_intermediate: bool = True


class OfficeRoundtripPipeline(TranslationPipeline):
    """
    Office Roundtrip Pipeline.
    
    Auto-detects the original Office format from PDF metadata and uses
    appropriate conversion tools.
    
    Workflow:
    1. Detect source format from PDF metadata
    2. Convert PDF to Office format (DOCX/PPTX/XLSX)
    3. Extract text from Office XML
    4. Generate translation file
    5. Parse translations and update XML
    6. Repack Office file
    7. Convert Office → PDF via LibreOffice
    """
    
    def __init__(self, config: OfficeRoundtripConfig):
        super().__init__(config)
        self.config: OfficeRoundtripConfig = config
        self._detected_format: Optional[SourceInfo] = None
        self._extraction_result: Optional[ExtractionResult] = None
        self._check_dependencies()
    
    def _check_dependencies(self) -> None:
        """Check available dependencies."""
        self._has_pdf2docx = False
        self._has_pptx = False
        self._has_openpyxl = False
        self._has_libreoffice = False
        
        try:
            import pdf2docx
            self._has_pdf2docx = True
        except ImportError:
            pass
        
        try:
            import pptx
            self._has_pptx = True
        except ImportError:
            pass
        
        try:
            import openpyxl
            self._has_openpyxl = True
        except ImportError:
            pass
        
        # Check LibreOffice
        lo_paths = [
            self.config.libreoffice_path,
            Path("C:/Program Files/LibreOffice/program/soffice.exe"),
            Path("C:/Program Files (x86)/LibreOffice/program/soffice.exe"),
            Path("/usr/bin/libreoffice"),
            Path("/usr/bin/soffice"),
            Path("/Applications/LibreOffice.app/Contents/MacOS/soffice"),
        ]
        for path in lo_paths:
            if path and path.exists():
                self.config.libreoffice_path = path
                self._has_libreoffice = True
                break
    
    @property
    def name(self) -> str:
        return "Office Roundtrip"
    
    @property
    def description(self) -> str:
        return "PDF -> Office (DOCX/PPTX/XLSX) -> Translate XML -> Office -> PDF"
    
    def derive_paths(self, input_path: Path) -> dict[str, Path]:
        """Derive paths including Office intermediates."""
        base = super().derive_paths(input_path)
        
        ext = self._get_office_extension()
        base.update({
            "office": input_path.with_suffix(ext),
            "office_translated": input_path.parent / f"{input_path.stem}_translated{ext}",
        })
        return base
    
    def _get_office_extension(self) -> str:
        """Get Office file extension based on detected/configured format."""
        if self.config.office_format == OfficeFormat.PPTX:
            return ".pptx"
        elif self.config.office_format == OfficeFormat.XLSX:
            return ".xlsx"
        else:
            return ".docx"
    
    def extract(self, input_path: Path) -> ExtractResult:
        """
        Extract text from PDF via Office format conversion.
        
        Steps:
        1. Detect source format if AUTO
        2. Convert PDF to Office format
        3. Extract XML from Office ZIP
        4. Parse text segments from XML
        5. Generate translation file
        """
        # Step 1: Detect source format if AUTO
        if self.config.office_format == OfficeFormat.AUTO:
            self._detected_format = detect_source_format(input_path)
            print(f"  Detected source: {self._detected_format.format.value} "
                  f"({self._detected_format.confidence:.0%} confidence)")
            
            # Map source format to office format
            if self._detected_format.format == SourceFormat.POWERPOINT:
                self.config.office_format = OfficeFormat.PPTX
            elif self._detected_format.format == SourceFormat.EXCEL:
                self.config.office_format = OfficeFormat.XLSX
            else:
                self.config.office_format = OfficeFormat.DOCX
        
        # Update paths with correct extension
        paths = self.derive_paths(input_path)
        
        # Step 2: Convert PDF to Office format
        pdf_size = input_path.stat().st_size / 1024  # KB
        print(f"  Input PDF: {pdf_size:.1f} KB")
        print(f"  Converting PDF -> {self.config.office_format.value.upper()}...")
        
        start_time = time.time()
        if self.config.office_format == OfficeFormat.DOCX:
            self._convert_pdf_to_docx(input_path, paths["office"])
        elif self.config.office_format == OfficeFormat.PPTX:
            self._convert_pdf_to_pptx(input_path, paths["office"])
        elif self.config.office_format == OfficeFormat.XLSX:
            self._convert_pdf_to_xlsx(input_path, paths["office"])
        
        elapsed = time.time() - start_time
        office_size = paths["office"].stat().st_size / 1024  # KB
        print(f"  Created {self.config.office_format.value.upper()}: {office_size:.1f} KB ({elapsed:.1f}s)")
        
        # Step 3: Extract XML and text segments
        print(f"  Extracting text from XML...")
        handler = get_handler(paths["office"])
        self._extraction_result = handler.extract()
        
        segments = self._extraction_result.segments
        print(f"  Found {len(segments)} text segments")
        
        # Step 4: Generate layout JSON
        layout = {
            "source_file": str(input_path),
            "pipeline": "office_roundtrip",
            "office_format": self.config.office_format.value,
            "office_path": str(paths["office"]),
            "blocks": [
                {
                    "block_id": seg.block_id,
                    "text": seg.text,
                    "xml_path": seg.xml_path,
                    "type": seg.metadata.get("type", "text"),
                }
                for seg in segments
            ],
            "block_order": [seg.block_id for seg in segments],
        }
        paths["layout"].write_text(
            json.dumps(layout, indent=2, ensure_ascii=False),
            encoding="utf-8"
        )
        
        # Step 5: Generate translation file
        lines = []
        for i, seg in enumerate(segments):
            text = seg.text.replace('\n', '\\n')
            lines.append(f"<{i}>{text}</{i}>")
        
        paths["translate"].write_text('\n'.join(lines), encoding="utf-8")
        paths["translated"].write_text('\n'.join(lines), encoding="utf-8")
        
        return ExtractResult(
            layout_path=paths["layout"],
            translate_path=paths["translate"],
            translated_template_path=paths["translated"],
            extra_files={
                "office": paths["office"],
                "format": self.config.office_format.value,
            },
        )
    
    def _convert_pdf_to_docx(self, pdf_path: Path, docx_path: Path) -> None:
        """Convert PDF to DOCX using pdf2docx."""
        if not self._has_pdf2docx:
            raise ImportError(
                "pdf2docx is required for PDF → DOCX conversion.\n"
                "Install with: pip install pdf2docx"
            )
        
        from pdf2docx import Converter
        
        cv = Converter(str(pdf_path))
        cv.convert(str(docx_path), start=0, end=None)
        cv.close()
    
    def _convert_pdf_to_pptx(self, pdf_path: Path, pptx_path: Path) -> None:
        """Convert PDF to PPTX by extracting text and creating slides."""
        if not self._has_pptx:
            raise ImportError(
                "python-pptx is required for PPTX creation.\n"
                "Install with: pip install python-pptx"
            )
        
        import fitz
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        
        # Extract text and layout from PDF
        doc = fitz.open(str(pdf_path))
        total_pages = len(doc)
        
        try:
            # Create presentation with same page size
            prs = Presentation()
            
            if total_pages > 0:
                first_page = doc[0]
                # Set slide size to match PDF (convert points to EMUs)
                prs.slide_width = int(first_page.rect.width * 914400 / 72)
                prs.slide_height = int(first_page.rect.height * 914400 / 72)
            
            blank_layout = prs.slide_layouts[6]  # Blank layout
            
            for page_num, page in enumerate(doc, 1):
                # Show progress
                sys.stdout.write(f"\r    Slide {page_num}/{total_pages}")
                sys.stdout.flush()
                
                slide = prs.slides.add_slide(blank_layout)
                blocks = page.get_text("dict")["blocks"]
                
                for block in blocks:
                    if block.get("type") == 0:  # Text block
                        bbox = block.get("bbox", [0, 0, 100, 20])
                        x, y, x1, y1 = bbox
                        
                        # Convert PDF points to inches
                        left = Inches(x / 72)
                        top = Inches(y / 72)
                        width = Inches((x1 - x) / 72)
                        height = Inches((y1 - y) / 72)
                        
                        # Collect text and get dominant font size
                        text_parts = []
                        font_size = 12
                        
                        for line in block.get("lines", []):
                            line_texts = []
                            for span in line.get("spans", []):
                                if span.get("text"):
                                    line_texts.append(span["text"])
                                    font_size = span.get("size", 12)
                            if line_texts:
                                text_parts.append(''.join(line_texts))
                        
                        if text_parts:
                            full_text = '\n'.join(text_parts)
                            
                            # Create text box
                            txBox = slide.shapes.add_textbox(left, top, width, height)
                            tf = txBox.text_frame
                            tf.word_wrap = True
                            
                            p = tf.paragraphs[0]
                            p.text = full_text
                            p.font.size = Pt(font_size)
            
            # Clear progress line
            sys.stdout.write(f"\r    Created {total_pages} slides    \n")
            sys.stdout.flush()
            
            prs.save(str(pptx_path))
            
        finally:
            doc.close()
    
    def _convert_pdf_to_xlsx(self, pdf_path: Path, xlsx_path: Path) -> None:
        """Convert PDF to XLSX by extracting tabular data."""
        if not self._has_openpyxl:
            raise ImportError(
                "openpyxl is required for XLSX creation.\n"
                "Install with: pip install openpyxl"
            )
        
        import fitz
        from openpyxl import Workbook
        from openpyxl.utils import get_column_letter
        
        doc = fitz.open(str(pdf_path))
        total_pages = len(doc)
        
        try:
            wb = Workbook()
            
            for page_num, page in enumerate(doc):
                # Show progress
                sys.stdout.write(f"\r    Sheet {page_num + 1}/{total_pages}")
                sys.stdout.flush()
                
                if page_num == 0:
                    ws = wb.active
                    ws.title = f"Sheet{page_num + 1}"
                else:
                    ws = wb.create_sheet(f"Sheet{page_num + 1}")
                
                # Extract text blocks and try to place them in cells
                blocks = page.get_text("dict")["blocks"]
                
                # Sort blocks by position (top to bottom, left to right)
                text_blocks = [b for b in blocks if b.get("type") == 0]
                text_blocks.sort(key=lambda b: (b["bbox"][1], b["bbox"][0]))
                
                # Group blocks into rows based on Y position
                rows = []
                current_row = []
                last_y = None
                
                for block in text_blocks:
                    y = block["bbox"][1]
                    if last_y is None or abs(y - last_y) < 15:  # Same row
                        current_row.append(block)
                    else:
                        if current_row:
                            rows.append(sorted(current_row, key=lambda b: b["bbox"][0]))
                        current_row = [block]
                    last_y = y
                
                if current_row:
                    rows.append(sorted(current_row, key=lambda b: b["bbox"][0]))
                
                # Write to worksheet
                for row_num, row_blocks in enumerate(rows, 1):
                    for col_num, block in enumerate(row_blocks, 1):
                        text_parts = []
                        for line in block.get("lines", []):
                            for span in line.get("spans", []):
                                if span.get("text"):
                                    text_parts.append(span["text"])
                        
                        cell_text = ' '.join(text_parts).strip()
                        if cell_text:
                            ws.cell(row=row_num, column=col_num, value=cell_text)
            
            # Clear progress line
            sys.stdout.write(f"\r    Created {total_pages} sheets    \n")
            sys.stdout.flush()
            
            wb.save(str(xlsx_path))
            
        finally:
            doc.close()
    
    def merge(
        self,
        input_path: Path,
        output_path: Path,
        translated_path: Path,
        layout_path: Path,
    ) -> MergeResult:
        """
        Apply translations and convert back to PDF.
        
        Steps:
        1. Load layout and parse translations
        2. Load Office file and update XML with translations
        3. Repack Office file
        4. Convert Office → PDF via LibreOffice
        """
        if not self._has_libreoffice:
            raise RuntimeError(
                "LibreOffice is required for Office -> PDF conversion.\n"
                "Install from: https://www.libreoffice.org/"
            )
        
        # Load layout
        layout = json.loads(layout_path.read_text(encoding="utf-8"))
        office_format = layout.get("office_format", "docx")
        office_path = Path(layout.get("office_path", ""))
        block_order = layout.get("block_order", [])
        
        if not office_path.exists():
            raise FileNotFoundError(f"Office file not found: {office_path}")
        
        # Parse translations
        translations = self._parse_translations(translated_path, block_order)
        print(f"  Parsed {len(translations)} translations")
        
        # Create translated Office file path
        ext = f".{office_format}"
        translated_office = input_path.parent / f"{input_path.stem}_translated{ext}"
        
        # Get handler and extract current state
        print(f"  Reading {office_format.upper()} structure...")
        handler = get_handler(office_path)
        extraction = handler.extract()
        
        # Update with translations
        print(f"  Updating {office_format.upper()} with {len(translations)} translations...")
        handler.update(translated_office, translations, extraction)
        
        # Get file size for progress info
        office_size = translated_office.stat().st_size / 1024  # KB
        print(f"  Office file: {office_size:.1f} KB")
        
        # Convert to PDF with progress spinner
        self._office_to_pdf(translated_office, output_path, office_format)
        
        # Keep intermediate files (original and translated Office files)
        print(f"  Kept: {office_path.name} ({office_format.upper()})")
        print(f"  Kept: {translated_office.name} ({office_format.upper()})")
        
        return MergeResult(
            output_path=output_path,
            blocks_processed=len(translations),
        )
    
    def _parse_translations(
        self,
        translated_path: Path,
        block_order: list[str],
    ) -> dict[str, str]:
        """Parse translated file."""
        content = translated_path.read_text(encoding="utf-8")
        translations = {}
        
        pattern = r'<(\d+)>(.*?)</\1>'
        for match in re.finditer(pattern, content, re.DOTALL):
            idx = int(match.group(1))
            text = match.group(2).replace('\\n', '\n')
            
            if idx < len(block_order):
                translations[block_order[idx]] = text
        
        return translations
    
    def _office_to_pdf(
        self,
        office_path: Path,
        pdf_path: Path,
        office_format: str = "docx",
    ) -> None:
        """Convert Office file to PDF via LibreOffice with progress indication."""
        cmd = [
            str(self.config.libreoffice_path),
            "--headless",
            "--convert-to", "pdf",
            "--outdir", str(pdf_path.parent),
            str(office_path),
        ]
        
        # Use spinner for progress during LibreOffice conversion
        with ProgressSpinner(f"Converting {office_format.upper()} to PDF via LibreOffice"):
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=300)
        
        if result.returncode != 0:
            raise RuntimeError(f"LibreOffice conversion failed: {result.stderr}")
        
        # LibreOffice names output based on input filename
        generated = office_path.with_suffix(".pdf")
        if generated.name != pdf_path.name:
            generated_in_outdir = pdf_path.parent / generated.name
            if generated_in_outdir.exists() and generated_in_outdir != pdf_path:
                shutil.move(str(generated_in_outdir), str(pdf_path))
    
    def get_translation_prompt(self, block_count: int) -> str:
        """Get translation prompt."""
        from ..translation_io import get_translation_prompt
        return get_translation_prompt(block_count, self.config.target_language)


def create_office_roundtrip_pipeline(
    target_language: str = "Hindi",
    office_format: OfficeFormat = OfficeFormat.AUTO,
    libreoffice_path: Optional[Path] = None,
    keep_intermediate: bool = True,
) -> OfficeRoundtripPipeline:
    """Factory function to create Office roundtrip pipeline."""
    config = OfficeRoundtripConfig(
        target_language=target_language,
        office_format=office_format,
        libreoffice_path=libreoffice_path,
        keep_intermediate=keep_intermediate,
    )
    return OfficeRoundtripPipeline(config)
